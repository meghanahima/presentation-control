from pptx import Presentation
from sentence_transformers import SentenceTransformer
from pptPath import ppt_path

def extract_slides_content(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
        slides.append(" ".join(content))
    
    return slides


slide_texts = extract_slides_content(ppt_path)
# Encode with SBERT
model = SentenceTransformer('all-MiniLM-L6-v2')
slide_embeddings = model.encode(slide_texts, convert_to_tensor=True)

import torch
import torch.nn.functional as F

def get_best_matching_slide(query):
    # Encode the user's query
    query_embedding = model.encode([query], convert_to_tensor=True)

    # Compute cosine similarity with each slide
    cosine_scores = F.cosine_similarity(query_embedding, slide_embeddings)

    # Get best matching slide index
    best_idx = torch.argmax(cosine_scores).item()
    best_score = cosine_scores[best_idx].item()

    return best_idx, best_score
